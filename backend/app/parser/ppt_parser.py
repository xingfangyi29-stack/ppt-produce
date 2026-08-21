import re
from pptx import Presentation
from typing import List, Dict, Any, Optional
from pathlib import Path

NUMBER_RE = re.compile(r"(?<!\d)(\d{1,3}(?:[\.,]\d+)?)(?:\s*%|\s*percent)?")
PERCENT_RE = re.compile(r"(\d{1,3}(?:[\.,]\d+)?)\s*%")
YEAR_RE = re.compile(r"\b(19|20)\d{2}\b")


def _text_from_shape(shape) -> str:
    texts = []
    if hasattr(shape, "text") and shape.text:
        texts.append(shape.text)
    elif hasattr(shape, "text_frame"):
        for p in shape.text_frame.paragraphs:
            runs = [r.text for r in p.runs if getattr(r, 'text', None)]
            if runs:
                texts.append(''.join(runs))
    return "\n".join(texts).strip()


def _extract_numbers_from_text(text: str) -> Dict[str, Any]:
    nums = []
    percents = []
    years = []
    for m in PERCENT_RE.finditer(text):
        try:
            val = float(m.group(1).replace(',', '.'))
            percents.append(val)
            nums.append(val)
        except Exception:
            continue
    for m in YEAR_RE.finditer(text):
        try:
            years.append(int(m.group(0)))
        except Exception:
            continue
    # generic numbers
    for m in NUMBER_RE.finditer(text):
        s = m.group(1)
        if PERCENT_RE.search(m.group(0)):
            continue
        try:
            val = float(s.replace(',', '.'))
            nums.append(val)
        except Exception:
            continue
    return {
        'numbers': nums,
        'percentages': percents,
        'years': years
    }


def _extract_tables_from_slide(slide) -> List[List[List[str]]]:
    tables = []
    for shape in slide.shapes:
        if shape.shape_type and getattr(shape, 'has_table', False):
            try:
                table = []
                for r in shape.table.rows:
                    row = []
                    for c in r.cells:
                        row.append(c.text.strip())
                    table.append(row)
                tables.append(table)
            except Exception:
                continue
    return tables


def _extract_texts_from_slide(slide) -> List[str]:
    texts = []
    for shape in slide.shapes:
        try:
            if getattr(shape, 'has_text_frame', False) and shape.text:
                texts.append(shape.text.strip())
        except Exception:
            continue
    return texts


def _is_question(line: str) -> bool:
    # heuristics: contains a question mark or starts with Q[0-9]
    if '?' in line:
        return True
    if re.match(r'^Q\s*\d+', line.strip(), re.IGNORECASE):
        return True
    # starts with common survey question words
    if re.match(r'^(How|What|Why|To what extent|To what degree)\b', line.strip(), re.IGNORECASE):
        return True
    return False


def parse_pptx(path: str) -> Dict[str, Any]:
    """
    Parse a PPTX file and extract structured survey-related information.

    Returns a dict with keys (some may be null):
      overall_score, participation_rate, benchmark, years, dimensions, drivers, slides
    Also returns a confidence dict and list of not_identified fields.
    """
    prs = Presentation(path)
    slides_out: List[Dict[str, Any]] = []

    overall_score = None
    participation_rate = None
    benchmark = None
    years_set = set()
    dimensions_set = set()
    drivers_set = set()

    confidence = {
        'overall_score': 0.0,
        'participation_rate': 0.0,
        'benchmark': 0.0,
        'years': 0.0,
        'dimensions': 0.0,
        'drivers': 0.0,
    }

    # iterate slides and collect heuristics
    for idx, slide in enumerate(prs.slides):
        slide_number = idx + 1
        title = None
        try:
            if slide.shapes.title and slide.shapes.title.text:
                title = slide.shapes.title.text.strip()
        except Exception:
            title = None

        texts = _extract_texts_from_slide(slide)
        full_text = '\n'.join(texts)
        tables = _extract_tables_from_slide(slide)

        # numbers in slide
        extracted = _extract_numbers_from_text(full_text)

        # detect questions
        questions = []
        question_candidates = []
        for t in texts:
            for line in t.split('\n'):
                line = line.strip()
                if not line:
                    continue
                if _is_question(line):
                    questions.append(line)
                else:
                    question_candidates.append(line)

        # detect dimension names: heuristics: short lines in ALL CAPS or Title Case near section titles
        dims = set()
        drivers = set()
        benchmarks = []
        part_rates = []
        overall_scores = []

        for t in texts:
            # split lines and inspect
            for line in [l.strip() for l in t.split('\n') if l.strip()]:
                low = line.lower()
                # participation rate
                if 'participat' in low or 'response rate' in low or 'response-rate' in low:
                    nums = _extract_numbers_from_text(line)
                    if nums['percentages']:
                        part_rates.extend(nums['percentages'])
                    elif nums['numbers']:
                        part_rates.extend(nums['numbers'])
                # overall engagement
                if 'overall engagement' in low or 'engagement score' in low or 'overall score' in low:
                    nums = _extract_numbers_from_text(line)
                    if nums['numbers'] or nums['percentages']:
                        overall_scores.extend(nums['numbers'] + nums['percentages'])
                # benchmark
                if 'benchmark' in low or 'external benchmark' in low or 'industry' in low:
                    nums = _extract_numbers_from_text(line)
                    if nums['numbers'] or nums['percentages']:
                        benchmarks.extend(nums['numbers'] + nums['percentages'])
                # drivers
                if 'driver' in low or 'key driver' in low:
                    drivers.add(line)
                # dimension heuristics: short lines, title-case and not ending with punctuation
                if len(line) <= 60 and (line[0].isupper() or line.isupper()) and len(line.split()) <= 6:
                    # avoid picking general sentences
                    if not line.endswith('.') and not line.endswith('?'):
                        dims.add(line)
                # collect years
                for y in YEAR_RE.findall(line):
                    try:
                        years_set.add(int(y))
                    except Exception:
                        pass

        # also inspect tables for numeric columns and header labels
        table_structs = []
        for table in tables:
            table_structs.append(table)
            # flatten headers and cells to search
            for row in table:
                for cell in row:
                    for y in YEAR_RE.findall(cell):
                        try:
                            years_set.add(int(y))
                        except Exception:
                            pass
                    low = cell.lower()
                    if 'participat' in low or 'response rate' in low:
                        nums = _extract_numbers_from_text(cell)
                        if nums['percentages']:
                            part_rates.extend(nums['percentages'])
                        elif nums['numbers']:
                            part_rates.extend(nums['numbers'])
                    if 'overall engagement' in low or 'engagement' in low:
                        nums = _extract_numbers_from_text(cell)
                        if nums['numbers'] or nums['percentages']:
                            overall_scores.extend(nums['numbers'] + nums['percentages'])
                    if 'benchmark' in low:
                        nums = _extract_numbers_from_text(cell)
                        if nums['numbers'] or nums['percentages']:
                            benchmarks.extend(nums['numbers'] + nums['percentages'])
                    if 'driver' in low:
                        drivers.add(cell)

        # aggregate findings
        if dims:
            dimensions_set.update(dims)
        if drivers:
            drivers_set.update(drivers)
        if benchmarks:
            benchmark_vals = [float(x) for x in benchmarks]
            # prefer first
        else:
            benchmark_vals = []

        # choose participation rate if found
        # prefer percentage values
        pr_val = None
        if part_rates:
            pr_val = float(part_rates[0])

        oscore_val = None
        if overall_scores:
            oscore_val = float(overall_scores[0])

        slide_record = {
            'slide_number': slide_number,
            'title': title,
            'text': full_text if full_text else None,
            'tables': table_structs if table_structs else None,
            'numbers': extracted['numbers'] if extracted['numbers'] else None,
            'percentages': extracted['percentages'] if extracted['percentages'] else None,
            'years': extracted['years'] if extracted['years'] else None,
            'questions': questions if questions else None,
            'dimension_candidates': list(dims) if dims else None,
            'driver_candidates': list(drivers) if drivers else None,
            'found_participation': pr_val,
            'found_overall_score': oscore_val,
            'found_benchmarks': benchmark_vals if benchmark_vals else None,
        }

        slides_out.append(slide_record)

        # update global picks: prefer identified values with labels
        if overall_score is None and oscore_val is not None:
            overall_score = oscore_val
            confidence['overall_score'] = 0.9
        if participation_rate is None and pr_val is not None:
            participation_rate = pr_val
            confidence['participation_rate'] = 0.9
        if benchmark is None and benchmark_vals:
            benchmark = float(benchmark_vals[0])
            confidence['benchmark'] = 0.9

    # Years global
    years = sorted(list(years_set)) if years_set else None
    if years:
        confidence['years'] = 0.8

    dimensions = list(dimensions_set) if dimensions_set else None
    if dimensions:
        confidence['dimensions'] = 0.7

    drivers = list(drivers_set) if drivers_set else None
    if drivers:
        confidence['drivers'] = 0.7

    # If some fields still None, attempt light-pass heuristic: search all slide texts for unlabeled numbers near keywords
    # (already mostly covered above). Do not invent values.

    not_identified = []
    for k in ['overall_score', 'participation_rate', 'benchmark', 'years', 'dimensions', 'drivers']:
        if locals()[k] is None:
            not_identified.append(k)

    result = {
        'overall_score': overall_score if overall_score is not None else None,
        'participation_rate': participation_rate if participation_rate is not None else None,
        'benchmark': benchmark if benchmark is not None else None,
        'years': years,
        'dimensions': dimensions,
        'drivers': drivers,
        'slides': slides_out,
        'confidence': confidence,
        'not_identified': not_identified,
        'ai_hooks': {
            'interpret_dimensions': None,
            'resolve_drivers': None,
            'notes': 'AI hooks present for future interpretation. No AI used in this pass.'
        }
    }

    return result
