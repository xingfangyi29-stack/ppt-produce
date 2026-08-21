from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.chart.data import CategoryChartData, ChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from typing import Dict, Any, List, Optional
from pathlib import Path
import math


def _safe_get(d: Dict, key: str):
    return d.get(key)


def _make_title_slide(prs: Presentation, overall: Optional[float], participation: Optional[float], company_name: Optional[str] = None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    left = Inches(0.5)
    top = Inches(0.5)
    width = prs.slide_width - Inches(1.0)

    # Header / brand area
    title_box = slide.shapes.add_textbox(left, top, width, Inches(1.0))
    tf = title_box.text_frame
    tf.text = company_name or "Management Discussion Deck"
    p = tf.paragraphs[0]
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = RGBColor(10, 25, 47)

    # Big numbers area
    num_left = Inches(0.75)
    num_top = Inches(1.8)
    num_width = Inches(4.5)
    num_height = Inches(2.5)

    # Overall engagement
    overall_box = slide.shapes.add_textbox(num_left, num_top, num_width, num_height)
    otf = overall_box.text_frame
    if overall is not None:
        otf.text = f"Overall Engagement\n{overall:.1f}"  # large number
    else:
        otf.text = "Overall Engagement\nNot identified"
    p0 = otf.paragraphs[0]
    p0.font.size = Pt(14)
    p0.font.bold = True
    p0.font.color.rgb = RGBColor(80, 86, 95)
    p1 = otf.paragraphs[1]
    p1.font.size = Pt(48)
    p1.font.bold = True
    p1.font.color.rgb = RGBColor(3, 37, 65)
    p1.alignment = PP_ALIGN.LEFT

    # Participation
    part_left = Inches(5.5)
    part_top = num_top
    part_width = Inches(3.0)
    part_height = Inches(2.5)
    part_box = slide.shapes.add_textbox(part_left, part_top, part_width, part_height)
    ptf = part_box.text_frame
    if participation is not None:
        ptf.text = f"Participation Rate\n{participation:.1f}%"
    else:
        ptf.text = "Participation Rate\nNot identified"
    p0 = ptf.paragraphs[0]
    p0.font.size = Pt(12)
    p0.font.color.rgb = RGBColor(80, 86, 95)
    p1 = ptf.paragraphs[1]
    p1.font.size = Pt(36)
    p1.font.bold = True
    p1.font.color.rgb = RGBColor(3, 37, 65)

    return slide


def _make_trend_slide(prs: Presentation, parsed: Dict[str, Any]):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    left = Inches(0.5)
    top = Inches(0.5)
    width = prs.slide_width - Inches(1.0)

    title_box = slide.shapes.add_textbox(left, top, width, Inches(0.6))
    tf = title_box.text_frame
    tf.text = "Trend & Benchmark"
    p = tf.paragraphs[0]
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = RGBColor(10, 25, 47)

    years = parsed.get('years') or []

    # Attempt to build an overall time series from slides tables if present
    # Look for slides with found_overall_score or table rows labeled 'Overall' — parser provides candidate values in slides
    series_years = []
    series_vals: List[float] = []
    for s in parsed.get('slides', []):
        if s.get('years') and s.get('found_overall_score') is not None:
            for y in s.get('years'):
                # if slide had multiple years but only one found_overall_score, we will map that value to last year
                pass
        # detect table rows named 'overall' inside parsed slides' tables
        if s.get('tables'):
            for table in s.get('tables'):
                # try header years
                header = None
                if len(table) >= 1:
                    # scan first two rows for years
                    for row in table[:2]:
                        ys = []
                        for cell in row:
                            if cell:
                                import re
                                m = re.search(r"(19|20)\\d{2}", cell)
                                if m:
                                    ys.append(int(m.group(0)))
                        if ys:
                            header = ys
                            break
                # search for a row label containing 'overall'
                for row in table:
                    if not row:
                        continue
                    label = (row[0] or '').lower()
                    if 'overall' in label:
                        # take numeric cells
                        vals = []
                        for cell in row[1:]:
                            if not cell:
                                vals.append(None)
                            else:
                                # extract number
                                import re
                                m = re.search(r"-?\d{1,3}(?:[\.,]\d+)?", cell)
                                if m:
                                    try:
                                        vals.append(float(m.group(0).replace(',', '.')))
                                    except Exception:
                                        vals.append(None)
                                else:
                                    vals.append(None)
                        if header:
                            for i, v in enumerate(vals[:len(header) if header else len(vals) ]):
                                if v is not None:
                                    series_years.append(header[i])
                                    series_vals.append(v)
                        else:
                            # map to global years if possible
                            for v in vals:
                                if v is not None:
                                    series_vals.append(v)
    # If we found series data, create a line chart
    chart_top = Inches(1.6)
    chart_left = Inches(0.7)
    chart_width = Inches(8.0)
    chart_height = Inches(3.8)

    chart_added = False
    if series_vals and series_years and len(series_years) == len(series_vals):
        data = CategoryChartData()
        categories = [str(y) for y in series_years]
        data.categories = categories
        data.add_series('Overall engagement', series_vals)
        # include benchmark if provided as a flat series
        benchmark = parsed.get('benchmark')
        if benchmark is not None:
            data.add_series('Benchmark', [benchmark for _ in series_vals])
        x, y, cx, cy = chart_left, chart_top, chart_width, chart_height
        chart = slide.shapes.add_chart(XL_CHART_TYPE.LINE, x, y, cx, cy, data).chart
        # style chart
        chart.has_legend = True
        chart_added = True
    else:
        # try simpler chart: single year comparing overall vs benchmark
        overall = parsed.get('overall_score')
        benchmark = parsed.get('benchmark')
        if overall is not None and benchmark is not None:
            data = CategoryChartData()
            data.categories = ['Overall', 'Benchmark']
            data.add_series('Score', (overall, benchmark))
            x, y, cx, cy = chart_left, chart_top, chart_width, chart_height
            chart = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, data).chart
            chart.has_legend = False
            chart_added = True

    note_top = Inches(5.1)
    note_box = slide.shapes.add_textbox(Inches(0.5), note_top, prs.slide_width - Inches(1.0), Inches(0.6))
    ntf = note_box.text_frame
    if chart_added:
        ntf.text = "Trend visualization (data-driven)."
    else:
        ntf.text = "Trend data not available in source report."
    ntf.paragraphs[0].font.size = Pt(12)
    ntf.paragraphs[0].font.color.rgb = RGBColor(90, 95, 102)

    return slide


def _make_strengths_slide(prs: Presentation, strengths: List[str], improvements: List[str]):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    left = Inches(0.5)
    top = Inches(0.5)
    width = prs.slide_width - Inches(1.0)

    title_box = slide.shapes.add_textbox(left, top, width, Inches(0.5))
    tf = title_box.text_frame
    tf.text = "What's Working & What Needs Attention"
    p = tf.paragraphs[0]
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = RGBColor(10, 25, 47)

    col_left = Inches(0.6)
    col_top = Inches(1.2)
    col_w = Inches(4.1)
    col_h = Inches(4.5)

    # Left: Strengths
    s_box = slide.shapes.add_textbox(col_left, col_top, col_w, col_h)
    stf = s_box.text_frame
    stf.text = "What's working"
    stf.paragraphs[0].font.size = Pt(14)
    stf.paragraphs[0].font.bold = True
    for s in strengths[:3]:
        p = stf.add_paragraph()
        p.level = 1
        p.text = f"• {s}"
        p.font.size = Pt(13)

    # Right: Improvements
    r_left = Inches(5.0)
    r_box = slide.shapes.add_textbox(r_left, col_top, col_w, col_h)
    rtf = r_box.text_frame
    rtf.text = "Needs attention"
    rtf.paragraphs[0].font.size = Pt(14)
    rtf.paragraphs[0].font.bold = True
    for im in improvements[:3]:
        p = rtf.add_paragraph()
        p.level = 1
        p.text = f"• {im}"
        p.font.size = Pt(13)

    return slide


def _make_insights_slide(prs: Presentation, insights: List[Dict[str, Any]]):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    left = Inches(0.5)
    top = Inches(0.5)
    width = prs.slide_width - Inches(1.0)

    title_box = slide.shapes.add_textbox(left, top, width, Inches(0.5))
    title_box.text_frame.text = "Key AI Insights"
    title_box.text_frame.paragraphs[0].font.size = Pt(18)
    title_box.text_frame.paragraphs[0].font.bold = True

    content_left = Inches(0.6)
    content_top = Inches(1.2)
    content_w = prs.slide_width - Inches(1.2)
    content_h = Inches(4.5)

    box = slide.shapes.add_textbox(content_left, content_top, content_w, content_h)
    tf = box.text_frame
    for idx, ins in enumerate(insights[:3]):
        title = ins.get('title') or f'Insight {idx+1}'
        what = ins.get('what_happened') or ''
        why = ins.get('why_it_matters') or ''
        # small paragraph per insight
        p = tf.add_paragraph()
        p.text = f"{title}: {what}"
        p.font.size = Pt(13)
        p.font.bold = True
        p2 = tf.add_paragraph()
        p2.text = why
        p2.level = 1
        p2.font.size = Pt(12)

    return slide


def _make_discussion_slide(prs: Presentation, discussion_areas: List[Dict[str, Any]]):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    left = Inches(0.5)
    top = Inches(0.5)
    width = prs.slide_width - Inches(1.0)

    title_box = slide.shapes.add_textbox(left, top, width, Inches(0.5))
    title_box.text_frame.text = "Discussion Areas"
    title_box.text_frame.paragraphs[0].font.size = Pt(18)
    title_box.text_frame.paragraphs[0].font.bold = True

    content_left = Inches(0.6)
    content_top = Inches(1.2)
    content_w = prs.slide_width - Inches(1.2)
    content_h = Inches(5.0)

    box = slide.shapes.add_textbox(content_left, content_top, content_w, content_h)
    tf = box.text_frame
    for idx, d in enumerate(discussion_areas[:2]):
        title = d.get('title') or f'Discussion {idx+1}'
        prompt = d.get('discussion_prompt') or ''
        supporting = d.get('supporting_data') or ''
        p = tf.add_paragraph()
        p.text = f"{title}"
        p.font.size = Pt(13)
        p.font.bold = True
        p2 = tf.add_paragraph()
        p2.text = prompt
        p2.level = 1
        p2.font.size = Pt(12)
        if supporting:
            p3 = tf.add_paragraph()
            p3.text = f"Evidence: {supporting}"
            p3.level = 2
            p3.font.size = Pt(11)
            p3.font.color.rgb = RGBColor(90, 95, 102)

    return slide


def _make_notes_slide(prs: Presentation, manager_notes: Optional[str]):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    left = Inches(0.5)
    top = Inches(0.5)
    width = prs.slide_width - Inches(1.0)

    title_box = slide.shapes.add_textbox(left, top, width, Inches(0.5))
    title_box.text_frame.text = "Manager Discussion Notes"
    title_box.text_frame.paragraphs[0].font.size = Pt(18)
    title_box.text_frame.paragraphs[0].font.bold = True

    box = slide.shapes.add_textbox(Inches(0.6), Inches(1.2), prs.slide_width - Inches(1.2), Inches(5.0))
    tf = box.text_frame
    if manager_notes:
        for line in str(manager_notes).split('\n'):
            p = tf.add_paragraph()
            p.text = line
            p.font.size = Pt(12)
    else:
        tf.text = "Manager notes go here."
        tf.paragraphs[0].font.size = Pt(12)

    return slide


def generate_deck(parsed_readonly: Dict[str, Any], edited: Dict[str, Any], output_path: str, template_path: Optional[str] = None) -> str:
    """
    Generate a 5-6 slide Management Discussion Deck PPTX.

    Inputs:
      - parsed_readonly: dict containing overall_score, participation_rate, benchmark, slides, years, etc. (read-only numeric anchors)
      - edited: dict containing edited text fields (insights, strengths, improvements, discussion_areas, manager_notes)
      - output_path: path to write the generated pptx
      - template_path: optional path to a .pptx template to preserve corporate branding; if None, a clean default is used

    Returns the path to the written PPTX file.
    """
    if template_path and Path(template_path).exists():
        prs = Presentation(template_path)
    else:
        prs = Presentation()

    # Ensure blank slide layout exists
    # Build slides
    overall = parsed_readonly.get('overall_score')
    participation = parsed_readonly.get('participation_rate')

    # Page 1
    _make_title_slide(prs, overall, participation, company_name=None)
    # Page 2: Trend + Benchmark
    _make_trend_slide(prs, parsed_readonly)
    # Page 3: What's Working & Needs Attention
    strengths = edited.get('strengths') or []
    improvements = edited.get('improvements') or []
    _make_strengths_slide(prs, strengths, improvements)
    # Page 4: Key AI Insights
    insights = edited.get('insights') or parsed_readonly.get('ai_insights') or []
    _make_insights_slide(prs, insights)
    # Page 5: Discussion Areas
    discussion_areas = edited.get('discussion_areas') or parsed_readonly.get('discussion_areas') or []
    _make_discussion_slide(prs, discussion_areas)
    # Page 6: Manager Notes
    manager_notes = edited.get('manager_notes') or ''
    _make_notes_slide(prs, manager_notes)

    out = Path(output_path)
    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out))
    return str(out)
